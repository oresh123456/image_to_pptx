"""
Module: extraction
==================
Opens a PPTX file and extracts the background image from each slide.

NotebookLM-exported PPTX files have exactly one large picture per slide that
covers the entire canvas. This module finds that picture and returns its raw
image bytes for downstream API calls.

Core functions (used in the pipeline):
  - extract_slide_inputs(pptx_path) -> (Presentation, list[dict]):
    Primary entry point used by pipeline.py. Opens the PPTX, iterates all
    slides, and for each slide that contains a picture, collects the python-pptx
    Slide object, the Picture shape, the raw image bytes, and the MIME type.
    The Presentation object is returned alongside so pipeline.py can call
    prs.save() at the end without reopening the file.

  - extract_slides(pptx_path) -> list[tuple[int, bytes, str]]:
    Simplified interface returning only (slide_idx, image_bytes, mime_type).
    Useful for ad-hoc scripts and unit testing.

Helper functions:
  - _iter_shapes(shapes): Recursively yields all leaf shapes, descending into
    group shapes. Group shapes are transparent containers in python-pptx.
  - _find_largest_picture(slide): Returns the picture shape with the largest
    pixel area, or None if the slide has no pictures.

Pipeline role: first stage, called once on the main thread before any parallel
  work begins. The returned Presentation object must stay alive on the main
  thread throughout the pipeline — python-pptx slide objects hold references
  into it and will break if it is garbage-collected.
"""

from __future__ import annotations

import logging

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

log = logging.getLogger(__name__)


def _iter_shapes(shapes):
    """
    Recursively yield all leaf shapes in a slide, descending into groups.

    python-pptx's MSO_SHAPE_TYPE.GROUP shapes are transparent containers;
    the actual content lives in their .shapes child collection. This generator
    flattens the tree so callers can treat all shapes uniformly.

    Args:
        shapes: A python-pptx SlideShapes collection or a GroupShapes collection.

    Yields:
        Individual non-group Shape objects in depth-first order.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _find_largest_picture(slide):
    """
    Return the picture shape with the largest bounding box area on a slide.

    In NotebookLM PPTX files the background image is always the largest picture.
    This heuristic also works for decks that have small decorative images
    alongside the main canvas image.

    Args:
        slide: A python-pptx Slide object.

    Returns:
        The Picture Shape with the greatest (width * height) product, or None
        if the slide contains no picture shapes at all.
    """
    pics = [
        s for s in _iter_shapes(slide.shapes)
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    if not pics:
        return None
    return max(pics, key=lambda p: (p.width or 0) * (p.height or 0))


def extract_slide_inputs(pptx_path: str) -> tuple:
    """
    Open a PPTX and collect all per-slide data needed by the pipeline.

    This is the primary extraction entry point. It returns the live
    Presentation object (which must be kept alive until prs.save() is
    called) along with a list of per-slide dicts containing everything
    pipeline.py needs — the Slide object for reconstruction, the Picture
    shape for position/size reference and replacement, and the raw image
    data for the API calls.

    Slides with no picture are logged and silently skipped.

    Args:
        pptx_path: Absolute or relative path to the input .pptx file.

    Returns:
        A 2-tuple (prs, slide_inputs) where:
          - prs is the open python-pptx Presentation object.
          - slide_inputs is a list of dicts, one per slide that has a picture:
            {
              "slide_idx":   int,    # 1-based slide number
              "slide":       Slide,  # python-pptx Slide object
              "pic":         Shape,  # the largest Picture shape on the slide
              "image_bytes": bytes,  # raw image data from pic.image.blob
              "mime_type":   str,    # e.g. "image/png" or "image/jpeg"
            }

    Raises:
        PackageNotFoundError: If pptx_path does not exist or is not a valid PPTX.
    """
    prs = Presentation(pptx_path)
    slide_inputs: list[dict] = []

    for idx, slide in enumerate(prs.slides, start=1):
        pic = _find_largest_picture(slide)
        if pic is None:
            log.info("Slide %d: no picture shape found — skipping.", idx)
            continue
        image = pic.image
        mime = image.content_type or "image/png"
        slide_inputs.append({
            "slide_idx":   idx,
            "slide":       slide,
            "pic":         pic,
            "image_bytes": image.blob,
            "mime_type":   mime,
        })
        log.debug(
            "Slide %d: found picture (%s, %d bytes).",
            idx, mime, len(image.blob),
        )

    return prs, slide_inputs


def extract_slides(pptx_path: str) -> list[tuple[int, bytes, str]]:
    """
    Simplified interface: open a PPTX and return image bytes per slide.

    Wraps extract_slide_inputs() and discards the Presentation object. Use
    this for quick scripts or testing; use extract_slide_inputs() for the
    full pipeline so the Presentation object stays alive.

    Args:
        pptx_path: Absolute or relative path to the input .pptx file.

    Returns:
        List of (slide_idx, image_bytes, mime_type) tuples, one per slide
        that contains at least one picture shape, in slide order.
    """
    _, slide_inputs = extract_slide_inputs(pptx_path)
    return [
        (si["slide_idx"], si["image_bytes"], si["mime_type"])
        for si in slide_inputs
    ]
