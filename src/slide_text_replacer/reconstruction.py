"""
Module: reconstruction
======================
Assembles the final PPTX slide: replaces the original rasterized image with
the inpainted (text-erased) version, then overlays one editable text box per
enriched region.

The result is a slide that looks visually identical to the original (background
restored by LaMa, text at the correct position and size) but where every text
element is a native PowerPoint text box — selectable, editable, copy-pasteable,
and searchable.

Core functions (used in the pipeline):
  - rebuild_slide(slide, old_pic, clean_image_bytes, enriched_regions, config)
    -> None:
    The single pipeline entry point. Captures the picture's EMU position and
    size before replacing it, then calls _replace_picture_image() and
    _add_text_box() for each region. Modifies the slide in place.

Helper functions:
  - _replace_picture_image(slide, old_pic, new_image_bytes) -> None:
    Inserts a new Picture shape at old_pic's position/size, then removes
    old_pic. python-pptx has no in-place image replacement API — this is the
    correct workaround (add then remove). Must be called before adding text
    boxes so text shapes render on top of the background.
  - _parse_color(color_hex) -> RGBColor:
    Parse a "#RRGGBB" hex string into a python-pptx RGBColor object.
  - _add_text_box(slide, region, pic_left_emu, pic_top_emu,
                  pic_width_emu, pic_height_emu, config) -> None:
    Create one text box for one region. Converts box_2d (0-1000 normalized)
    to EMU using the picture's reference frame. Sets RTL, font, size, color.

Pipeline role: final per-slide stage, called on the MAIN thread after all
  ThreadPoolExecutor futures have completed. python-pptx is not thread-safe;
  all mutations to Slide objects happen here under the main thread.

Font sizing (notes.md §6.3):
  pt = px * config.font_px_to_pt  (default 0.75, the 96 DPI convention)
  Clamped to [config.font_min_pt, config.font_max_pt] (default 8–60 pt).

RTL rendering (notes.md §6.1 and §6.2):
  Every paragraph gets alignment=RIGHT + set_paragraph_rtl().
  Every run gets set_run_font() which sets both <a:latin> and <a:cs>.
"""

from __future__ import annotations

import copy
import io
import logging

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Emu, Pt

from slide_text_replacer.config import Config
from slide_text_replacer.pptx_helpers import set_run_font, set_paragraph_rtl
from slide_text_replacer.schemas import EnrichedRegion

log = logging.getLogger(__name__)


def _replace_picture_image(slide, old_pic, new_image_bytes: bytes) -> None:
    """
    Replace a slide picture's image while keeping its position and size.

    python-pptx does not expose an in-place image replacement API, so we
    add a new Picture shape at the same (left, top, width, height) and then
    remove the old shape from the XML tree. The new shape lands at the top
    of the z-order; call this before adding text boxes so text renders above
    the background.

    Args:
        slide:           The python-pptx Slide that contains old_pic.
        old_pic:         The existing Picture shape to replace.
        new_image_bytes: Raw bytes of the replacement image (PNG or JPEG).

    Returns:
        None. Modifies slide.shapes in place.
    """
    left, top     = old_pic.left, old_pic.top
    width, height = old_pic.width, old_pic.height

    slide.shapes.add_picture(
        io.BytesIO(new_image_bytes), left, top, width=width, height=height
    )
    old_pic._element.getparent().remove(old_pic._element)


def _parse_color(color_hex: str) -> RGBColor:
    """
    Parse a "#RRGGBB" hex color string into a python-pptx RGBColor.

    Args:
        color_hex: Color string in "#RRGGBB" format (the leading # is optional).

    Returns:
        RGBColor instance.

    Raises:
        ValueError: If color_hex is not a valid 6-digit hex color.
    """
    h = color_hex.lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Expected 6-digit hex color, got: {color_hex!r}")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _add_text_box(
    slide,
    region: EnrichedRegion,
    pic_left_emu: int,
    pic_top_emu: int,
    pic_width_emu: int,
    pic_height_emu: int,
    config: Config,
) -> None:
    """
    Add a single editable text box to the slide for one enriched region.

    Converts the region's normalized 0-1000 bounding box to EMU coordinates
    using the picture's EMU position and size as the reference frame. Creates
    the text box with no margins, word wrap enabled, top-anchored, with a
    right-aligned RTL Hebrew paragraph containing a single run styled with
    the region's font, size, color, and bold setting.

    Args:
        slide:          The python-pptx Slide to add the text box to.
        region:         The enriched region with all visual metadata.
        pic_left_emu:   Left edge of the slide picture in EMU.
        pic_top_emu:    Top edge of the slide picture in EMU.
        pic_width_emu:  Width of the slide picture in EMU.
        pic_height_emu: Height of the slide picture in EMU.
        config:         Pipeline config for font sizing parameters.

    Returns:
        None. Adds a TextBox shape to the slide in place.
    """
    ymin, xmin, ymax, xmax = region.box_2d

    # Denormalize from 0-1000 to EMU using the picture's reference frame.
    left_emu = pic_left_emu + int(xmin / 1000 * pic_width_emu)
    top_emu  = pic_top_emu  + int(ymin / 1000 * pic_height_emu)
    w_emu    = max(Emu(10_000), int((xmax - xmin) / 1000 * pic_width_emu))
    h_emu    = max(Emu(10_000), int((ymax - ymin) / 1000 * pic_height_emu))

    tb = slide.shapes.add_textbox(left_emu, top_emu, w_emu, h_emu)
    tf = tb.text_frame
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap       = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size       = MSO_AUTO_SIZE.NONE

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(para)

    run = para.add_run()
    run.text = region.text

    # Convert px → pt and clamp (notes.md §6.3).
    font_pt = region.font_size_px * config.font_px_to_pt
    font_pt = max(config.font_min_pt, min(config.font_max_pt, round(font_pt)))
    run.font.size = Pt(font_pt)
    run.font.bold = (region.font_weight == "bold")

    try:
        run.font.color.rgb = _parse_color(region.color)
    except ValueError:
        log.debug(
            "Unparseable color %r for %r — using black.", region.color, region.text[:20]
        )
        run.font.color.rgb = RGBColor(0, 0, 0)

    # Set font for both Latin and complex-script (Hebrew) — see notes.md §6.1.
    set_run_font(run, region.font_family)

    # Remove the text box border so it's invisible.
    tb.line.fill.background()


def rebuild_slide(
    slide,
    old_pic,
    clean_image_bytes: bytes,
    enriched_regions: list[EnrichedRegion],
    config: Config,
) -> None:
    """
    Replace a slide's background image and add editable text overlays.

    This is the main reconstruction entry point. It:
      1. Reads and stores the picture's current EMU position and size.
      2. Calls _replace_picture_image() to swap in the inpainted background.
      3. Calls _add_text_box() for each enriched region to create the overlay.

    Errors in individual text box creation are caught and logged without
    aborting the slide — a partial overlay is better than a failed slide.

    The slide is modified in place. Call prs.save() after all slides are done.

    Args:
        slide:             The python-pptx Slide object to modify.
        old_pic:           The original Picture shape (used for position/size
                           reference, then removed by _replace_picture_image).
        clean_image_bytes: Inpainted image bytes to use as the new background.
        enriched_regions:  Text regions with visual metadata for overlay.
        config:            Pipeline config containing font sizing constants.

    Returns:
        None. Modifies the slide in place.
    """
    # Capture position before the shape is removed.
    pic_left   = old_pic.left
    pic_top    = old_pic.top
    pic_width  = old_pic.width
    pic_height = old_pic.height

    _replace_picture_image(slide, old_pic, clean_image_bytes)
    log.debug(
        "Picture replaced. Adding %d text box(es).", len(enriched_regions)
    )

    for region in enriched_regions:
        try:
            _add_text_box(
                slide, region,
                pic_left, pic_top, pic_width, pic_height,
                config,
            )
        except Exception as exc:
            log.warning(
                "Failed to add text box for %r: %s", region.text[:30], exc
            )


def add_candidate_slide(
    prs,
    clean_image_bytes: bytes,
    enriched_regions: list[EnrichedRegion],
    original_slide,
    config: Config,
):
    """
    Create a new blank slide with the inpainted image and text overlays.

    Uses the blank slide layout and positions the image at the same bounds
    as the original slide's picture.

    Args:
        prs:                The Presentation object.
        clean_image_bytes:  Inpainted image bytes (PNG).
        enriched_regions:   Text regions with visual metadata.
        original_slide:     The original slide (used for slide dimensions reference).
        config:             Pipeline config.

    Returns:
        The newly created slide object.
    """
    # Use blank layout (index 6) or fall back to first available
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(layout)

    # Use full slide dimensions for the image
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add the clean image covering the full slide
    pic = new_slide.shapes.add_picture(
        io.BytesIO(clean_image_bytes),
        Emu(0), Emu(0),
        width=slide_width, height=slide_height,
    )

    # Add text boxes
    for region in enriched_regions:
        try:
            _add_text_box(
                new_slide, region,
                Emu(0), Emu(0), slide_width, slide_height,
                config,
            )
        except Exception as exc:
            log.warning(
                "Failed to add text box for %r on candidate slide: %s",
                region.text[:30], exc,
            )

    return new_slide


def interleave_candidate_slides(prs, slide_inputs: list[dict], added_slides: list) -> None:
    """
    Reorder slides so candidates are interleaved: [1a, 1b, 2a, 2b, ...].

    After reconstruction, original slides are in their original positions and
    candidate slides are appended at the end. This function moves them into
    interleaved order.

    Args:
        prs:           The Presentation object.
        slide_inputs:  Original slide input dicts (in order).
        added_slides:  List of (original_slide_idx, new_slide) tuples.
    """
    if not added_slides:
        return

    # Build a map: original_slide_idx → new_slide
    idx_to_new_slide = {idx: slide for idx, slide in added_slides}

    # Build desired order: for each original slide, insert its candidate after it
    slide_list = prs.slides._sldIdLst
    # Get current slide elements in order
    sldId_elements = list(slide_list)

    # Map slide objects to their sldId elements by rId
    slide_id_map = {}
    for sldId_elem in sldId_elements:
        rId = sldId_elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        slide_id_map[rId] = sldId_elem

    # Build new order: for each original slide_input, place original then candidate
    new_order = []
    original_rIds = set()

    for si in slide_inputs:
        idx = si["slide_idx"]
        # Find the rId for the original slide
        orig_slide = si["slide"]
        orig_rId = None
        for rel in prs.part.rels.values():
            if hasattr(rel, '_target') and rel._target is orig_slide.part:
                orig_rId = rel.rId
                break
        if orig_rId is None:
            # Fallback: find by slide part
            for rId, sldId_elem in slide_id_map.items():
                # We'll use position-based matching below
                pass

        # Add original slide's sldId
        if orig_rId and orig_rId in slide_id_map:
            new_order.append(slide_id_map[orig_rId])
            original_rIds.add(orig_rId)

        # Add candidate slide's sldId if it exists
        if idx in idx_to_new_slide:
            cand_slide = idx_to_new_slide[idx]
            cand_rId = None
            for rel in prs.part.rels.values():
                if hasattr(rel, '_target') and rel._target is cand_slide.part:
                    cand_rId = rel.rId
                    break
            if cand_rId and cand_rId in slide_id_map:
                new_order.append(slide_id_map[cand_rId])
                original_rIds.add(cand_rId)

    # Add any remaining slides not in our map (shouldn't happen, but safe)
    for sldId_elem in sldId_elements:
        rId = sldId_elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId not in original_rIds:
            new_order.append(sldId_elem)

    # Clear and re-add in new order
    for elem in list(slide_list):
        slide_list.remove(elem)
    for elem in new_order:
        slide_list.append(elem)
