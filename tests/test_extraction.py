"""
Tests for extraction.py — PPTX slide image extraction.

Verifies I/O contracts documented in docs/modules/extraction.md.
All tests are local — no API calls. python-pptx Presentations created in-memory.
"""

import io

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from slide_text_replacer.extraction import (
    _iter_shapes,
    _find_largest_picture,
    extract_slide_inputs,
    extract_slides,
)


# ── Helpers ──────────────────────────────────────────────────────────────────

def _make_image_bytes(w: int = 100, h: int = 100) -> bytes:
    """Create a minimal PNG image of the given dimensions."""
    img = Image.new("RGB", (w, h), (128, 128, 128))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_pptx_with_pictures(tmp_path, slides_config: list[list[tuple[int, int]]]) -> str:
    """Create a PPTX with the given slide configuration, return path."""
    prs = Presentation()
    for pics in slides_config:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if not pics:
            slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1_000_000), Emu(500_000))
        for w, h in pics:
            img_bytes = _make_image_bytes(w // 9144, h // 9144 or 1)
            slide.shapes.add_picture(
                io.BytesIO(img_bytes), Emu(0), Emu(0), width=Emu(w), height=Emu(h)
            )
    path = str(tmp_path / "test.pptx")
    prs.save(path)
    return path


# ── _iter_shapes: output contract ────────────────────────────────────────────

def test_iter_shapes_yields_leaf_shapes():
    """Input: slide with textbox + picture → Output: yields both as leaf shapes."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1_000_000), Emu(500_000))
    slide.shapes.add_picture(
        io.BytesIO(_make_image_bytes()), Emu(0), Emu(0), width=Emu(500_000), height=Emu(500_000)
    )
    shapes = list(_iter_shapes(slide.shapes))
    assert len(shapes) == 2


# ── _find_largest_picture: output contract ───────────────────────────────────

def test_find_largest_picture_returns_biggest():
    """Input: slide with 2 pictures → Output: the one with larger area."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        io.BytesIO(_make_image_bytes(10, 10)),
        Emu(0), Emu(0), width=Emu(100_000), height=Emu(100_000),
    )
    slide.shapes.add_picture(
        io.BytesIO(_make_image_bytes(200, 200)),
        Emu(0), Emu(0), width=Emu(9_000_000), height=Emu(6_000_000),
    )
    pic = _find_largest_picture(slide)
    assert pic is not None
    assert pic.width == Emu(9_000_000)


def test_find_largest_picture_returns_none_when_no_pics():
    """Input: slide with only textboxes → Output: None."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1_000_000), Emu(500_000))
    pic = _find_largest_picture(slide)
    assert pic is None


# ── extract_slide_inputs(): output contract ──────────────────────────────────

def test_extract_slide_inputs_correct_count(tmp_path):
    """Input: 3-slide PPTX with pictures → Output: 3 slide input dicts."""
    path = _make_pptx_with_pictures(tmp_path, [
        [(9_144_000, 6_858_000)],
        [(9_144_000, 6_858_000)],
        [(9_144_000, 6_858_000)],
    ])
    prs, inputs = extract_slide_inputs(path)
    assert len(inputs) == 3


def test_extract_slide_inputs_dict_keys(tmp_path):
    """Input: PPTX with 1 picture slide → Output: dict with documented keys."""
    path = _make_pptx_with_pictures(tmp_path, [[(9_144_000, 6_858_000)]])
    _, inputs = extract_slide_inputs(path)
    expected_keys = {"slide_idx", "slide", "pic", "image_bytes", "mime_type"}
    assert set(inputs[0].keys()) == expected_keys


def test_extract_slide_inputs_skips_no_picture_slides(tmp_path):
    """Input: slides 1,3 have pictures, slide 2 has none → Output: 2 dicts with correct slide_idx."""
    path = _make_pptx_with_pictures(tmp_path, [
        [(9_144_000, 6_858_000)],  # slide 1: has picture
        [],                         # slide 2: no picture
        [(9_144_000, 6_858_000)],  # slide 3: has picture
    ])
    _, inputs = extract_slide_inputs(path)
    assert len(inputs) == 2
    assert inputs[0]["slide_idx"] == 1
    assert inputs[1]["slide_idx"] == 3


# ── extract_slides(): output contract ────────────────────────────────────────

def test_extract_slides_returns_tuples(tmp_path):
    """Input: PPTX with 1 picture → Output: list of (int, bytes, str) tuples."""
    path = _make_pptx_with_pictures(tmp_path, [[(9_144_000, 6_858_000)]])
    result = extract_slides(path)
    assert len(result) == 1
    idx, img_bytes, mime = result[0]
    assert isinstance(idx, int)
    assert isinstance(img_bytes, bytes)
    assert isinstance(mime, str)


def test_extract_slides_empty_pptx(tmp_path):
    """Input: PPTX with no pictures → Output: empty list."""
    path = _make_pptx_with_pictures(tmp_path, [[], []])
    result = extract_slides(path)
    assert result == []
