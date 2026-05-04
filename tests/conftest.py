"""
Shared pytest fixtures used across multiple test files.

No API calls — all fixtures produce synthetic in-memory objects.
Integration fixtures (real_config, test_pptx_path, output_dir) skip
when prerequisites are missing.
"""

import io
import json
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from slide_text_replacer.config import Config, load_config
from slide_text_replacer.schemas import Region, EnrichedRegion


_FIXTURES_DIR = Path(__file__).parent / "fixtures"
_TEST_PPTX = _FIXTURES_DIR / "test_input.pptx"
_OUTPUT_DIR = Path(__file__).parent / "output"


@pytest.fixture()
def sample_image_bytes() -> bytes:
    """Minimal 200x200 grey PNG image."""
    img = Image.new("RGB", (200, 200), (200, 200, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture()
def sample_regions() -> list[Region]:
    """Three Region objects matching tests/fixtures/sample_ocr.json."""
    raw = json.loads((_FIXTURES_DIR / "sample_ocr.json").read_text(encoding="utf-8"))
    return [
        Region(
            text=item["text"],
            box_2d=tuple(item["box_2d"]),
            font_size_px=float(item["font_size_px"]),
        )
        for item in raw
    ]


@pytest.fixture()
def sample_enriched_regions() -> list[EnrichedRegion]:
    """Three EnrichedRegion objects matching tests/fixtures/sample_enriched.json."""
    raw = json.loads((_FIXTURES_DIR / "sample_enriched.json").read_text(encoding="utf-8"))
    return [
        EnrichedRegion(
            text=item["text"],
            box_2d=tuple(item["box_2d"]),
            font_size_px=float(item["font_size_px"]),
            font_family=item["font_family"],
            font_weight=item["font_weight"],
            color=item["color"],
        )
        for item in raw
    ]


@pytest.fixture()
def dummy_config() -> Config:
    """Config with fake API keys and default settings."""
    return Config(
        gemini_api_key="fake-gemini-key",
        replicate_token="fake-replicate-token",
    )


@pytest.fixture()
def pptx_with_picture(sample_image_bytes: bytes, tmp_path: Path) -> Path:
    """Create a minimal PPTX with one slide containing a Picture shape."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        io.BytesIO(sample_image_bytes),
        Emu(0), Emu(0),
        width=Emu(9_144_000), height=Emu(6_858_000),
    )
    pptx_path = tmp_path / "test_input.pptx"
    prs.save(str(pptx_path))
    return pptx_path


# ---------------------------------------------------------------------------
# Integration fixtures (skip when prerequisites missing)
# ---------------------------------------------------------------------------


@pytest.fixture()
def real_config() -> Config:
    """Load real config.toml. Skip if keys missing."""
    try:
        return load_config()
    except (RuntimeError, FileNotFoundError):
        pytest.skip("config.toml missing or API keys not set")


@pytest.fixture()
def test_pptx_path() -> Path:
    """Resolve tests/fixtures/test_input.pptx. Skip if missing."""
    if not _TEST_PPTX.exists():
        pytest.skip(f"{_TEST_PPTX} not found — place a real PPTX there")
    return _TEST_PPTX


@pytest.fixture()
def output_dir() -> Path:
    """Ensure tests/output/ exists and return its Path."""
    _OUTPUT_DIR.mkdir(exist_ok=True)
    return _OUTPUT_DIR
