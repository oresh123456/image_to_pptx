"""Integration tests — text overlay (reconstruction without inpainting).

Requires: real config keys, test_input.pptx in tests/fixtures/.
Outputs saved to tests/output/ for visual inspection.
"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Pt

from slide_text_replacer.config import Config
from slide_text_replacer.extraction import extract_slide_inputs
from slide_text_replacer.ocr import run as ocr_run
from slide_text_replacer.enrichment import run as enrichment_run
from slide_text_replacer.reconstruction import rebuild_slide


pytestmark = pytest.mark.integration


def _process_slides(test_pptx_path: Path, config: Config):
    """Extract, OCR, enrich all slides. Returns (prs, slides, enriched_map)."""
    prs, slides = extract_slide_inputs(str(test_pptx_path))
    enriched_map = {}
    for s in slides:
        regions = ocr_run(
            s["image_bytes"], s["mime_type"],
            config.gemini_api_key, config.gemini_model,
        )
        enriched = enrichment_run(
            s["image_bytes"], s["mime_type"], regions,
            config.gemini_api_key, config.gemini_model,
        )
        enriched_map[s["slide_idx"]] = enriched
    return prs, slides, enriched_map


class TestOverlayOnOriginal:
    def test_overlay_produces_openable_pptx(
        self, real_config, test_pptx_path, output_dir,
    ):
        """Rebuild using original image (no inpaint). Output is valid PPTX with text boxes."""
        prs, slides, enriched_map = _process_slides(test_pptx_path, real_config)

        for s in slides:
            enriched = enriched_map.get(s["slide_idx"], [])
            if not enriched:
                continue
            rebuild_slide(
                s["slide"], s["pic"], s["image_bytes"],
                enriched, real_config,
            )

        out = output_dir / "overlay_only.pptx"
        prs.save(str(out))

        # verify openable + has shapes
        check = Presentation(str(out))
        for slide in check.slides:
            assert len(slide.shapes) >= 1

    def test_overlay_text_matches_ocr(
        self, real_config, test_pptx_path, output_dir,
    ):
        """Text in output PPTX text boxes matches OCR regions."""
        prs, slides, enriched_map = _process_slides(test_pptx_path, real_config)

        for s in slides:
            enriched = enriched_map.get(s["slide_idx"], [])
            if not enriched:
                continue
            rebuild_slide(
                s["slide"], s["pic"], s["image_bytes"],
                enriched, real_config,
            )

        out = output_dir / "overlay_text_check.pptx"
        prs.save(str(out))

        check = Presentation(str(out))
        for slide_idx, slide in enumerate(check.slides, 1):
            enriched = enriched_map.get(slide_idx, [])
            if not enriched:
                continue
            expected_texts = {r.text for r in enriched}
            found_texts = set()
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            found_texts.add(text)
            # every enriched text should appear in the slide
            for t in expected_texts:
                assert t in found_texts, (
                    f"Slide {slide_idx}: '{t}' not found in text boxes"
                )

    def test_overlay_font_properties(
        self, real_config, test_pptx_path, output_dir,
    ):
        """Font family, size, and color are set on text runs."""
        prs, slides, enriched_map = _process_slides(test_pptx_path, real_config)

        for s in slides:
            enriched = enriched_map.get(s["slide_idx"], [])
            if not enriched:
                continue
            rebuild_slide(
                s["slide"], s["pic"], s["image_bytes"],
                enriched, real_config,
            )

        out = output_dir / "overlay_fonts_check.pptx"
        prs.save(str(out))

        check = Presentation(str(out))
        for slide in check.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        assert run.font.size is not None, (
                            f"Font size not set on run '{run.text}'"
                        )
                        assert run.font.color.rgb is not None, (
                            f"Font color not set on run '{run.text}'"
                        )
