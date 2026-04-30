"""
Tests for pptx_helpers.py — XML helpers for Hebrew font and RTL rendering.

Verifies I/O contracts documented in docs/modules/pptx_helpers.md.
All tests are local — no API calls. Verifies lxml element structure directly.
"""

import pytest
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

from slide_text_replacer.pptx_helpers import set_run_font, set_paragraph_rtl


# ── Helpers ──────────────────────────────────────────────────────────────────

def _make_run_and_para():
    """Create a slide with a textbox, return (paragraph, run) for testing."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(3_000_000), Emu(500_000))
    para = txBox.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "שלום עולם"
    return para, run


# ── set_run_font(): output contract ──────────────────────────────────────────
# Input: (run, font_name) → Output: run XML has <a:latin>, <a:cs>, lang="he-IL".

def test_set_run_font_adds_latin_element():
    """Input: font_name="Heebo" → Output: <a:latin typeface="Heebo"> in run XML."""
    _, run = _make_run_and_para()
    set_run_font(run, "Heebo")
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    assert latin is not None
    assert latin.get("typeface") == "Heebo"


def test_set_run_font_adds_cs_element():
    """Input: font_name="Heebo" → Output: <a:cs typeface="Heebo"> (required for Hebrew)."""
    _, run = _make_run_and_para()
    set_run_font(run, "Heebo")
    rPr = run._r.get_or_add_rPr()
    cs = rPr.find(qn("a:cs"))
    assert cs is not None
    assert cs.get("typeface") == "Heebo"


def test_set_run_font_sets_lang_he_il():
    """Input: any font → Output: lang="he-IL" on run properties."""
    _, run = _make_run_and_para()
    set_run_font(run, "Rubik")
    rPr = run._r.get_or_add_rPr()
    assert rPr.get("lang") == "he-IL"


def test_set_run_font_works_with_all_palette_fonts():
    """Input: each of the 5 palette fonts → Output: <a:cs typeface=font> for each."""
    fonts = ["Heebo", "Rubik", "Assistant", "Frank Ruhl Libre", "Heebo Black"]
    for font in fonts:
        _, run = _make_run_and_para()
        set_run_font(run, font)
        rPr = run._r.get_or_add_rPr()
        cs = rPr.find(qn("a:cs"))
        assert cs is not None and cs.get("typeface") == font


def test_set_run_font_clears_existing_font_elements():
    """Input: called twice with different fonts → Output: no duplicate elements, last font wins."""
    _, run = _make_run_and_para()
    set_run_font(run, "Heebo")
    set_run_font(run, "Rubik")
    rPr = run._r.get_or_add_rPr()
    latin_elements = rPr.findall(qn("a:latin"))
    cs_elements = rPr.findall(qn("a:cs"))
    assert len(latin_elements) == 1
    assert len(cs_elements) == 1
    assert latin_elements[0].get("typeface") == "Rubik"
    assert cs_elements[0].get("typeface") == "Rubik"


# ── set_paragraph_rtl(): output contract ─────────────────────────────────────
# Input: paragraph → Output: rtl="1" on <a:pPr>.

def test_set_paragraph_rtl_sets_attribute():
    """Input: paragraph → Output: rtl="1" on paragraph properties."""
    para, _ = _make_run_and_para()
    set_paragraph_rtl(para)
    pPr = para._p.get_or_add_pPr()
    assert pPr.get("rtl") == "1"


def test_set_paragraph_rtl_is_idempotent():
    """Input: called twice → Output: still rtl="1" (no duplication)."""
    para, _ = _make_run_and_para()
    set_paragraph_rtl(para)
    set_paragraph_rtl(para)
    pPr = para._p.get_or_add_pPr()
    assert pPr.get("rtl") == "1"


def test_set_paragraph_rtl_does_not_affect_other_paragraphs():
    """Input: RTL set on para1 → Output: para2 in same frame is not affected."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(3_000_000), Emu(1_000_000))
    tf = txBox.text_frame
    para1 = tf.paragraphs[0]
    para1.add_run().text = "first paragraph"
    para2 = tf.add_paragraph()
    para2.add_run().text = "second paragraph"

    set_paragraph_rtl(para1)

    pPr1 = para1._p.get_or_add_pPr()
    pPr2 = para2._p.get_or_add_pPr()
    assert pPr1.get("rtl") == "1"
    assert pPr2.get("rtl") != "1"
