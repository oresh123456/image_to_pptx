"""Integration tests — full pipeline end-to-end.

Requires: real config keys, test_input.pptx in tests/fixtures/.
Outputs saved to tests/output/ for visual inspection.
"""

from pathlib import Path

import pytest
from pptx import Presentation

from slide_text_replacer.pipeline import run_pipeline


pytestmark = pytest.mark.integration


class TestFullPipeline:
    def test_produces_output(
        self, real_config, test_pptx_path, output_dir,
    ):
        """Pipeline produces a valid PPTX with same slide count."""
        out = output_dir / "full_result.pptx"
        run_pipeline(str(test_pptx_path), str(out), real_config)

        assert out.exists(), "Output file not created"
        orig = Presentation(str(test_pptx_path))
        result = Presentation(str(out))
        assert len(result.slides) == len(orig.slides), (
            f"Slide count mismatch: {len(result.slides)} vs {len(orig.slides)}"
        )

    def test_slides_have_text_boxes(
        self, real_config, test_pptx_path, output_dir,
    ):
        """Each slide in output has at least one text box."""
        out = output_dir / "full_result_textboxes.pptx"
        run_pipeline(str(test_pptx_path), str(out), real_config)

        result = Presentation(str(out))
        for i, slide in enumerate(result.slides, 1):
            has_text = any(s.has_text_frame for s in slide.shapes)
            assert has_text, f"Slide {i}: no text boxes found"
