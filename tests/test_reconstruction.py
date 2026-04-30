"""
Tests for reconstruction.py — slide rebuilding with inpainted images and text overlays.

Verifies I/O contracts documented in docs/modules/reconstruction.md.
All tests are local — no API calls. Real python-pptx objects used.
"""

import io

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from slide_text_replacer.config import Config
from slide_text_replacer.reconstruction import _parse_color, _add_text_box, rebuild_slide
from slide_text_replacer.schemas import EnrichedRegion


# ── Helpers ──────────────────────────────────────────────────────────────────

def _make_image_bytes(w: int = 200, h: int = 200) -> bytes:
    img = Image.new("RGB", (w, h), (200, 200, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_slide_with_picture():
    """Create a Presentation with one slide + Picture, return (slide, pic)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(
        io.BytesIO(_make_image_bytes()),
        Emu(0), Emu(0), width=Emu(9_144_000), height=Emu(6_858_000),
    )
    return slide, pic


def _make_config() -> Config:
    return Config(gemini_api_key="k", replicate_token="t")


def _make_region(**overrides) -> EnrichedRegion:
    defaults = dict(
        text="שלום עולם", box_2d=(100, 100, 200, 500), font_size_px=24.0,
        font_family="Heebo", font_weight="regular", color="#000000",
    )
    defaults.update(overrides)
    return EnrichedRegion(**defaults)


# ── _parse_color(): output contract ──────────────────────────────────────────
# Input: "#RRGGBB" or "RRGGBB" string → Output: RGBColor.

def test_parse_color_valid_hex():
    """Input: "#1a3a8a" → Output: RGBColor(0x1a, 0x3a, 0x8a)."""
    color = _parse_color("#1a3a8a")
    assert color == RGBColor(0x1A, 0x3A, 0x8A)


def test_parse_color_without_hash():
    """Input: "FF0000" (no hash) → Output: RGBColor(255, 0, 0)."""
    color = _parse_color("FF0000")
    assert color == RGBColor(255, 0, 0)


# ── _parse_color(): error handling ───────────────────────────────────────────

def test_parse_color_invalid_length_raises():
    """Input: "#FFF" (3-digit) → raises ValueError."""
    with pytest.raises(ValueError, match="6-digit"):
        _parse_color("#FFF")


def test_parse_color_invalid_hex_raises():
    """Input: "#GGHHII" → raises ValueError."""
    with pytest.raises(ValueError):
        _parse_color("#GGHHII")


# ── _add_text_box(): output contract ─────────────────────────────────────────
# Input: slide, EnrichedRegion, pic bounds, config → Output: new shape on slide.

def test_add_text_box_creates_shape():
    """Input: region → Output: slide.shapes count increases by 1."""
    slide, pic = _make_slide_with_picture()
    initial_count = len(slide.shapes)
    _add_text_box(slide, _make_region(), pic.left, pic.top, pic.width, pic.height, _make_config())
    assert len(slide.shapes) == initial_count + 1


def test_add_text_box_sets_text():
    """Input: region.text = "כותרת ראשית" → Output: run text matches."""
    slide, pic = _make_slide_with_picture()
    _add_text_box(slide, _make_region(text="כותרת ראשית"), pic.left, pic.top, pic.width, pic.height, _make_config())
    tb = slide.shapes[-1]
    run_text = tb.text_frame.paragraphs[0].runs[0].text
    assert run_text == "כותרת ראשית"


def test_add_text_box_clamps_font_size_max():
    """Input: font_size_px=200 (→ 150pt) → Output: clamped to 60pt."""
    slide, pic = _make_slide_with_picture()
    _add_text_box(slide, _make_region(font_size_px=200.0), pic.left, pic.top, pic.width, pic.height, _make_config())
    tb = slide.shapes[-1]
    font_size = tb.text_frame.paragraphs[0].runs[0].font.size
    assert font_size == Pt(60)


def test_add_text_box_clamps_font_size_min():
    """Input: font_size_px=2 (→ 1.5pt) → Output: clamped to 8pt."""
    slide, pic = _make_slide_with_picture()
    _add_text_box(slide, _make_region(font_size_px=2.0), pic.left, pic.top, pic.width, pic.height, _make_config())
    tb = slide.shapes[-1]
    font_size = tb.text_frame.paragraphs[0].runs[0].font.size
    assert font_size == Pt(8)


def test_add_text_box_bold_weight():
    """Input: font_weight="bold" → Output: run.font.bold is True."""
    slide, pic = _make_slide_with_picture()
    _add_text_box(slide, _make_region(font_weight="bold"), pic.left, pic.top, pic.width, pic.height, _make_config())
    tb = slide.shapes[-1]
    assert tb.text_frame.paragraphs[0].runs[0].font.bold is True


# ── rebuild_slide(): output contract ─────────────────────────────────────────
# Input: slide, old_pic, clean_bytes, regions, config → Output: slide modified in place.

def test_rebuild_slide_adds_text_boxes():
    """Input: 3 regions → Output: 3 textbox shapes added to slide."""
    slide, pic = _make_slide_with_picture()
    clean_image = _make_image_bytes()
    regions = [_make_region(text="one"), _make_region(text="two"), _make_region(text="three")]

    rebuild_slide(slide, pic, clean_image, regions, _make_config())

    textbox_count = sum(
        1 for s in slide.shapes
        if s.has_text_frame and s.text_frame.paragraphs[0].runs
    )
    assert textbox_count == 3
